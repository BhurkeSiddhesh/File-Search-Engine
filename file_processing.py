import os
from docx import Document
from pypdf import PdfReader
from pptx import Presentation
from openpyxl import load_workbook

def extract_text(filepath):
    """
    Extracts text from a file based on its extension.
    """
    ext = os.path.splitext(filepath)[1].lower()

    try:
        if ext == '.txt':
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        elif ext == '.docx':
            doc = Document(filepath)
            return "\n".join([para.text for para in doc.paragraphs])
        elif ext == '.pdf':
            with open(filepath, 'rb') as f:
                reader = PdfReader(f)
                return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
        elif ext == '.pptx':
            prs = Presentation(filepath)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        elif ext == '.xlsx':
            workbook = load_workbook(filepath, read_only=True)
            text = []
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value:
                            text.append(str(cell.value))
            return "\n".join(text)
        else:
            print(f"Unsupported file type: {ext}")
            return None
    except Exception as e:
        print(f"Error extracting text from {filepath}: {e}")
        return None
